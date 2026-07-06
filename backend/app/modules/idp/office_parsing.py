"""Text extraction for office/text/email formats (Phase 1 universal ingestion).

Each function takes raw bytes and returns the document's plain text, ready to
feed into ``search_tsv`` exactly like the PDF/OCR paths. None of these formats
have a meaningful page count, so callers treat them as a single logical page.
"""

import io
from email import message_from_bytes, policy
from email.message import Message
from html.parser import HTMLParser


class _HtmlTextStripper(HTMLParser):
    """Minimal HTML-to-text fallback for email bodies with no text/plain part.

    Stdlib-only — good enough for "make it searchable", not a rendering engine.
    """

    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []

    def handle_data(self, data: str) -> None:
        self._chunks.append(data)

    def text(self) -> str:
        return " ".join(self._chunks)


def _strip_html(html: str) -> str:
    stripper = _HtmlTextStripper()
    stripper.feed(html)
    return stripper.text()


def extract_docx_text(data: bytes) -> str:
    """Concatenate paragraph and table cell text from a .docx file."""
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts)


def extract_xlsx_text(data: bytes) -> str:
    """Concatenate every cell value across every sheet of an .xlsx file."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        parts.append(str(cell.value))
        return "\n".join(parts)
    finally:
        wb.close()


def extract_pptx_text(data: bytes) -> str:
    """Concatenate every text frame across every slide of a .pptx file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def extract_plain_text(data: bytes) -> str:
    """Decode txt/csv/md bytes (UTF-8, falling back to Latin-1)."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


def _email_body_text(msg: Message) -> str:
    """Prefer the text/plain part; fall back to a stripped text/html part."""
    if msg.is_multipart():
        plain_parts: list[str] = []
        html_parts: list[str] = []
        for part in msg.walk():
            if part.is_multipart():
                continue
            content_type = part.get_content_type()
            try:
                payload = part.get_content()
            except Exception:
                continue
            if not isinstance(payload, str):
                continue
            if content_type == "text/plain":
                plain_parts.append(payload)
            elif content_type == "text/html":
                html_parts.append(payload)
        if plain_parts:
            return "\n".join(plain_parts)
        if html_parts:
            return "\n".join(_strip_html(p) for p in html_parts)
        return ""

    try:
        payload = msg.get_content()
    except Exception:
        return ""
    if not isinstance(payload, str):
        return ""
    if msg.get_content_type() == "text/html":
        return _strip_html(payload)
    return payload


def extract_email_text(data: bytes) -> str:
    """Extract headers (From/To/Subject/Date) + body text from an .eml file."""
    msg = message_from_bytes(data, policy=policy.default)
    header_lines = [
        f"{key}: {msg.get(key)}"
        for key in ("From", "To", "Cc", "Subject", "Date")
        if msg.get(key)
    ]
    body = _email_body_text(msg)
    return "\n".join([*header_lines, "", body])
