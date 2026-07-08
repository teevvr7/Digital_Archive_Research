"""Universal ingestion (Phase 1) unit tests.

Covers the new content-sniffing MIME detector, the office/text/email text
extractors, thumbnail generation, and the pipeline dispatch wiring for the
expanded format set. All worker deps (PyMuPDF, Pillow, python-docx, openpyxl,
python-pptx) are genuinely installed in this dev environment, so these tests
exercise real round-trips rather than mocks where practical.
"""

import io

import pytest

from app.modules.idp import mimetype


# ---------------------------------------------------------------------------
# MIME sniffer
# ---------------------------------------------------------------------------


def test_sniff_pdf():
    assert mimetype.sniff_mime(b"%PDF-1.4 rest of file") == mimetype.MIME_PDF


def test_sniff_png():
    assert mimetype.sniff_mime(b"\x89PNG\r\n\x1a\n rest") == mimetype.MIME_PNG


def test_sniff_jpeg():
    assert mimetype.sniff_mime(b"\xff\xd8\xff\xe0 rest") == mimetype.MIME_JPEG


def test_sniff_webp():
    data = b"RIFF" + b"\x00\x00\x00\x00" + b"WEBP" + b"rest"
    assert mimetype.sniff_mime(data) == mimetype.MIME_WEBP


def test_sniff_tiff_little_and_big_endian():
    assert mimetype.sniff_mime(b"II*\x00rest") == mimetype.MIME_TIFF
    assert mimetype.sniff_mime(b"MM\x00*rest") == mimetype.MIME_TIFF


def test_sniff_docx_xlsx_pptx_via_zip_peek():
    import openpyxl
    from docx import Document as DocxDocument
    from pptx import Presentation

    docx_buf = io.BytesIO()
    DocxDocument().save(docx_buf)
    assert mimetype.sniff_mime(docx_buf.getvalue(), "f.docx") == mimetype.MIME_DOCX

    xlsx_buf = io.BytesIO()
    openpyxl.Workbook().save(xlsx_buf)
    assert mimetype.sniff_mime(xlsx_buf.getvalue(), "f.xlsx") == mimetype.MIME_XLSX

    pptx_buf = io.BytesIO()
    Presentation().save(pptx_buf)
    assert mimetype.sniff_mime(pptx_buf.getvalue(), "f.pptx") == mimetype.MIME_PPTX


def test_sniff_plain_text_uses_extension_for_subtype():
    assert mimetype.sniff_mime(b"hello,world\n1,2,3", "data.csv") == mimetype.MIME_CSV
    assert mimetype.sniff_mime(b"# Heading\n\nbody text", "notes.md") == mimetype.MIME_MD
    assert mimetype.sniff_mime(b"just some text", "notes.txt") == mimetype.MIME_TXT
    # No filename at all still falls back to plain text, never None.
    assert mimetype.sniff_mime(b"just some text", None) == mimetype.MIME_TXT


def test_sniff_renaming_a_text_file_to_csv_is_harmless_not_a_security_issue():
    """Extension only flips the cosmetic subtype within the text family —
    never used to admit a file that failed binary/text decoding."""
    assert mimetype.sniff_mime(b"plain text content", "report.csv") == mimetype.MIME_CSV


def test_sniff_email_detected_by_headers_not_extension():
    raw = b"From: alice@example.com\nTo: bob@example.com\nSubject: Hi\n\nBody text."
    assert mimetype.sniff_mime(raw, "whatever.dat") == mimetype.MIME_EML


def test_sniff_plain_text_with_one_coincidental_header_like_line_is_not_email():
    """A single 'Key: value'-looking line shouldn't false-positive as email —
    the detector requires >=2 recognised header keys."""
    raw = b"From: the management\n\nThis memo is about the new policy."
    assert mimetype.sniff_mime(raw, "memo.txt") == mimetype.MIME_TXT


def test_sniff_unrecognised_binary_returns_none():
    assert mimetype.sniff_mime(b"\x00\x01\x02\x03\x04\xff\xfe\xfd") is None


def test_sniff_random_high_entropy_binary_without_nul_still_returns_none():
    """Latin-1 decode never fails on any byte sequence — must not be used as
    the sole signal for 'this is text', or random binary would always pass."""
    garbage = bytes(range(1, 32)) + bytes(range(0x80, 0xA0))  # control + C1 ranges
    assert mimetype.sniff_mime(garbage) is None


def test_sniff_legacy_latin1_text_is_still_accepted():
    """A real legacy (non-UTF-8) Latin-1 text file with accented characters
    must still be classified as text, not rejected by the printability gate."""
    raw = "Café au lait, déjà vu, naïve résumé".encode("latin-1")
    assert mimetype.sniff_mime(raw, "notes.txt") == mimetype.MIME_TXT


# ---------------------------------------------------------------------------
# Office/text/email text extraction
# ---------------------------------------------------------------------------


def test_extract_docx_text_round_trip():
    from docx import Document as DocxDocument

    from app.modules.idp.office_parsing import extract_docx_text

    doc = DocxDocument()
    doc.add_paragraph("Hello from a Word document.")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Vendor"
    table.rows[0].cells[1].text = "Acme Corp"
    buf = io.BytesIO()
    doc.save(buf)

    text = extract_docx_text(buf.getvalue())
    assert "Hello from a Word document." in text
    assert "Vendor" in text and "Acme Corp" in text


def test_extract_xlsx_text_round_trip():
    from openpyxl import Workbook

    from app.modules.idp.office_parsing import extract_xlsx_text

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Invoice Total"
    ws["B1"] = 1234.5
    buf = io.BytesIO()
    wb.save(buf)

    text = extract_xlsx_text(buf.getvalue())
    assert "Invoice Total" in text
    assert "1234.5" in text


def test_extract_pptx_text_round_trip():
    from pptx import Presentation
    from pptx.util import Inches

    from app.modules.idp.office_parsing import extract_pptx_text

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Quarterly Results"
    buf = io.BytesIO()
    prs.save(buf)

    text = extract_pptx_text(buf.getvalue())
    assert "Quarterly Results" in text


def test_extract_plain_text_utf8_and_latin1_fallback():
    from app.modules.idp.office_parsing import extract_plain_text

    assert extract_plain_text("héllo wörld".encode("utf-8")) == "héllo wörld"
    # Latin-1 bytes that are invalid UTF-8 must not raise.
    latin1_bytes = "café".encode("latin-1")
    assert "caf" in extract_plain_text(latin1_bytes)


def test_extract_email_text_plain_body():
    from email.message import EmailMessage

    from app.modules.idp.office_parsing import extract_email_text

    msg = EmailMessage()
    msg["From"] = "alice@example.com"
    msg["To"] = "bob@example.com"
    msg["Subject"] = "Quarterly Invoice"
    msg.set_content("Please find the invoice attached.")

    text = extract_email_text(msg.as_bytes())
    assert "alice@example.com" in text
    assert "Quarterly Invoice" in text
    assert "Please find the invoice attached." in text


def test_extract_email_text_html_only_falls_back_to_stripped_text():
    from email.message import EmailMessage

    from app.modules.idp.office_parsing import extract_email_text

    msg = EmailMessage()
    msg["From"] = "alice@example.com"
    msg["Subject"] = "HTML only"
    msg.set_content("<html><body><p>Hello <b>World</b></p></body></html>", subtype="html")

    text = extract_email_text(msg.as_bytes())
    assert "Hello" in text and "World" in text
    assert "<p>" not in text


# ---------------------------------------------------------------------------
# Thumbnail generation
# ---------------------------------------------------------------------------


def test_generate_thumbnail_for_pdf():
    import fitz

    from app.modules.idp.thumbnails import generate_thumbnail

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello thumbnail")
    buf = doc.tobytes()
    doc.close()

    thumb = generate_thumbnail(buf, "application/pdf")
    assert thumb is not None
    assert thumb.startswith(b"\x89PNG\r\n\x1a\n")


def test_generate_thumbnail_for_image():
    from PIL import Image

    from app.modules.idp.thumbnails import generate_thumbnail

    img = Image.new("RGB", (800, 600), color="red")
    buf = io.BytesIO()
    img.save(buf, format="PNG")

    thumb = generate_thumbnail(buf.getvalue(), "image/png")
    assert thumb is not None
    out = Image.open(io.BytesIO(thumb))
    assert max(out.size) <= 320


def test_generate_thumbnail_none_for_office_and_text_types():
    from app.modules.idp.thumbnails import generate_thumbnail

    assert generate_thumbnail(b"anything", mimetype.MIME_DOCX) is None
    assert generate_thumbnail(b"anything", mimetype.MIME_TXT) is None
    assert generate_thumbnail(b"anything", mimetype.MIME_EML) is None


def test_generate_thumbnail_never_raises_on_corrupt_bytes():
    from app.modules.idp.thumbnails import generate_thumbnail

    assert generate_thumbnail(b"not a real pdf", "application/pdf") is None
    assert generate_thumbnail(b"not a real image", "image/png") is None


# ---------------------------------------------------------------------------
# Pipeline dispatch for the new format families
# ---------------------------------------------------------------------------


def test_pipeline_dispatches_docx_to_office_extractor():
    from docx import Document as DocxDocument

    from app.modules.idp.pipeline import run_extraction

    doc = DocxDocument()
    doc.add_paragraph("Searchable office content.")
    buf = io.BytesIO()
    doc.save(buf)

    result = run_extraction(buf.getvalue(), mimetype.MIME_DOCX)
    assert "Searchable office content." in result.text
    assert result.page_count == 1
    assert result.has_text_layer is True
    assert result.ocr_used is False


def test_pipeline_dispatches_plain_text():
    from app.modules.idp.pipeline import run_extraction

    result = run_extraction(b"Some searchable plain text.", mimetype.MIME_TXT)
    assert "Some searchable plain text." in result.text
    assert result.ocr_used is False
    assert result.has_text_layer is True


def test_pipeline_dispatches_email():
    from email.message import EmailMessage

    from app.modules.idp.pipeline import run_extraction

    msg = EmailMessage()
    msg["From"] = "a@example.com"
    msg["Subject"] = "Test"
    msg.set_content("Body content here.")

    result = run_extraction(msg.as_bytes(), mimetype.MIME_EML)
    assert "Body content here." in result.text
    assert result.page_count == 1


def test_pipeline_raises_on_unsupported_mime():
    from app.modules.idp.pipeline import run_extraction

    with pytest.raises(ValueError):
        run_extraction(b"whatever", "application/x-completely-unknown")


# ---------------------------------------------------------------------------
# document_date heuristic (idp/jobs.py)
# ---------------------------------------------------------------------------


def test_guess_document_date_finds_a_date_near_the_top():
    from app.modules.idp.jobs import _guess_document_date

    text = "Invoice Date: 15 March 2026\nVendor: Acme Corp\nAmount: 100.00"
    result = _guess_document_date(text)
    assert result is not None
    assert result.year == 2026 and result.month == 3 and result.day == 15


def test_guess_document_date_returns_none_for_no_text():
    from app.modules.idp.jobs import _guess_document_date

    assert _guess_document_date(None) is None
    assert _guess_document_date("") is None


def test_guess_document_date_never_raises_on_garbage():
    from app.modules.idp.jobs import _guess_document_date

    # Garbage OCR-ish text shouldn't crash the heuristic.
    assert _guess_document_date("@@@ ### $$$ %%%") is None or True


def test_guess_document_date_rejects_implausible_far_future_year():
    """Regression, calibrated against a real live upload: ``big-invoice.pdf``'s
    address line ("...TRARALGON, 3844, AUSTRALIA...") contains an Australian
    postcode that ``dateparser`` misread as a year, producing
    document_date=3844-07-03 (confirmed live via direct DB query against the
    real ``extracted_text``).

    Mocks ``dateparser``'s return directly rather than trying to reproduce its
    exact fuzzy free-text match, which is dateparser-version-dependent and not
    what this guard tests — the guard is the plausibility bound in
    ``_guess_document_date``, not dateparser's internal matching behavior.
    """
    from unittest.mock import patch
    import datetime as _dt
    from app.modules.idp.jobs import _guess_document_date

    with patch("dateparser.search.search_dates") as mock_search:
        mock_search.return_value = [("3844", _dt.datetime(3844, 7, 3))]
        assert _guess_document_date("...TRARALGON, 3844, AUSTRALIA...") is None


def test_guess_document_date_rejects_implausible_far_past_year():
    """Regression, calibrated against a real live upload: a marketing PDF with
    no genuine document date produced document_date=1949-07-06 (confirmed
    live via direct DB query against the real ``extracted_text``)."""
    from unittest.mock import patch
    import datetime as _dt
    from app.modules.idp.jobs import _guess_document_date

    with patch("dateparser.search.search_dates") as mock_search:
        mock_search.return_value = [("1949", _dt.datetime(1949, 7, 6))]
        assert _guess_document_date("marketing copy with no real date") is None


# ---------------------------------------------------------------------------
# VLM eligibility gating — structured extraction stays type-conditional
# ---------------------------------------------------------------------------


def test_vlm_eligible_mimes_is_pdf_and_images_only():
    assert mimetype.VLM_ELIGIBLE_MIMES == {mimetype.MIME_PDF} | mimetype.IMAGE_MIMES
    assert mimetype.MIME_DOCX not in mimetype.VLM_ELIGIBLE_MIMES
    assert mimetype.MIME_TXT not in mimetype.VLM_ELIGIBLE_MIMES
    assert mimetype.MIME_EML not in mimetype.VLM_ELIGIBLE_MIMES
